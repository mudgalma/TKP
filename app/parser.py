"""Document parsing utilities.

Converts uploaded files (PDF, PPTX, DOCX, TXT) to plain Markdown text
so downstream AI components can consume them uniformly.
"""

from __future__ import annotations

import asyncio
import io
import logging
import os
import tempfile
from pptx import Presentation
from docx import Document

logger = logging.getLogger(__name__)

# Word heading style → Markdown ATX heading prefix
_HEADING_MARKERS: dict[str, str] = {
    "Heading 1": "#",
    "Heading 2": "##",
    "Heading 3": "###",
}


# ── PDF ──────────────────────────────────────────────────────────────────────

async def parse_pdf(content: bytes) -> str:
    """Convert PDF bytes to Markdown text using LlamaParse."""
    from llama_parse import LlamaParse
    import tempfile
    import os

    # Ensure API key is available
    api_key = os.environ.get("LLAMA_CLOUD_API_KEY")
    if not api_key:
        raise ValueError("LLAMA_CLOUD_API_KEY environment variable is not set. Please set it in your .env file.")

    # Initialize the parser
    parser = LlamaParse(
        api_key=api_key,
        result_type="markdown",
        verbose=True
    )

    # LlamaParse expects a file path, so we write the bytes to a temporary file
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_pdf:
        temp_pdf.write(content)
        temp_pdf_path = temp_pdf.name

    try:
        # Asynchronously parse the document
        documents = await parser.aload_data(temp_pdf_path)
        
        # Combine the text from all parsed pages/chunks
        markdown_text = "\n\n".join([doc.text for doc in documents])
        return markdown_text
    finally:
        # Clean up the temporary file
        if os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)


# ── PPTX ─────────────────────────────────────────────────────────────────────

def parse_pptx(content: bytes) -> str:
    """Convert PPTX bytes to Markdown.

    Each slide becomes a level-2 heading; all text frames within
    that slide are appended below it.
    """

    prs = Presentation(io.BytesIO(content))
    lines: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Slide {slide_num}")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    lines.append(text)
        lines.append("")  # blank line between slides

    return "\n".join(lines)


# ── DOCX ─────────────────────────────────────────────────────────────────────

def parse_docx(content: bytes) -> str:
    """Convert DOCX bytes to Markdown.

    Word heading styles (Heading 1/2/3) are mapped to ATX Markdown
    headings (#/##/###). All other paragraphs are emitted as plain text.
    """


    doc = Document(io.BytesIO(content))
    lines: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Check if paragraph uses a heading style
        for heading_style, marker in _HEADING_MARKERS.items():
            if heading_style in para.style.name:
                text = f"{marker} {text}"
                break

        lines.append(text)

    return "\n\n".join(lines)


# ── DISPATCHER ───────────────────────────────────────────────────────────────

async def parse_file(filename: str, content_type: str, content: bytes) -> str:
    """Dispatch to the correct parser based on MIME type.

    Args:
        filename:     Original file name (used only for logging).
        content_type: MIME type reported by the client.
        content:      Raw file bytes.

    Returns:
        Extracted text as a Markdown string.

    Raises:
        ValueError: If the content type is not supported.
    """
    logger.debug(
        "Parsing '%s' (content_type=%s, size=%d bytes)",
        filename, content_type, len(content),
    )

    if content_type == "application/pdf":
        return await parse_pdf(content)

    if "pptx" in content_type or "powerpoint" in content_type:
        return parse_pptx(content)

    if "docx" in content_type or "wordprocessingml" in content_type:
        return parse_docx(content)

    if content_type == "text/plain":
        return content.decode("utf-8", errors="ignore")

    raise ValueError(
        f"Unsupported content type '{content_type}' for file '{filename}'. "
        "Accepted: application/pdf, pptx, docx, text/plain."
    )
