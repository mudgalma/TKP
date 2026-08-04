"""Document parsing utilities.

Converts uploaded files (PDF, PPTX, DOCX, TXT) into a list of
(page_number, text_block) tuples so downstream chunking preserves exact
page provenance for citations.

Routing strategy (cost-aware):
    "mostly_text"      → PyMuPDF (free, instant, CPU-only)
    "tables" / "equations" / "scanned" → LlamaParse (cloud, higher quality)
    "auto"             → heuristic: inspect embedded image count / text density

Failure contract:
    All public functions raise ParsingError on failure — raw fitz / llama
    exceptions must not leak past this module boundary.
"""

from __future__ import annotations

import asyncio
import io
import logging
import os
import tempfile
from typing import Literal

from tenacity import (
    RetryError,
    retry,
    retry_if_exception_type,
    stop_after_attempt,
    wait_exponential,
)

from app.config import settings
from app.exceptions import ParsingError

logger = logging.getLogger(__name__)

# Accepted hint values from the user / routing logic
DocumentHint = Literal["mostly_text", "tables", "equations", "scanned", "auto"]

# --------------------------------------------------------------------------- #
#  Page-block type                                                             #
# --------------------------------------------------------------------------- #

PageBlock = tuple[int, str]  # (1-indexed page number, text content)


# --------------------------------------------------------------------------- #
#  PyMuPDF path (free, local, text-heavy docs)                                #
# --------------------------------------------------------------------------- #

def _parse_pdf_pymupdf(content: bytes) -> list[PageBlock]:
    """Extract text from a PDF page-by-page using PyMuPDF (fitz)."""
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise ParsingError("PyMuPDF (fitz) is not installed.") from exc

    try:
        doc = fitz.open(stream=content, filetype="pdf")
        blocks: list[PageBlock] = []

        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                blocks.append((page_num, text))

        doc.close()
        logger.info("PyMuPDF: extracted %d text blocks from PDF.", len(blocks))
        return blocks
    except Exception as exc:
        raise ParsingError(f"PyMuPDF failed to parse PDF: {exc}") from exc


# --------------------------------------------------------------------------- #
#  LlamaParse path (cloud, tables / equations / scanned OCR)                  #
# --------------------------------------------------------------------------- #

@retry(
    retry=retry_if_exception_type((TimeoutError, ConnectionError)),
    stop=stop_after_attempt(settings.llm_max_retries),
    wait=wait_exponential(multiplier=1, min=2, max=30),
    reraise=True,
)
async def _parse_pdf_llamaparse(content: bytes, filename: str = "document.pdf") -> list[PageBlock]:
    """Extract text from a PDF using LlamaParse (cloud API)."""
    try:
        from llama_parse import LlamaParse
    except ImportError as exc:
        raise ParsingError("llama-parse is not installed.") from exc

    api_key = settings.llama_cloud_api_key
    if not api_key:
        raise ParsingError(
            "LLAMA_CLOUD_API_KEY is not set. Cannot use LlamaParse routing."
        )

    parser = LlamaParse(
        api_key=api_key,
        result_type="markdown",
        verbose=False,
    )

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        documents = await parser.aload_data(tmp_path)
    except RetryError as exc:
        raise ParsingError("LlamaParse failed after retries (network/timeout).") from exc
    except Exception as exc:
        # Check HTTP status — don't retry on 4xx (except 429 already handled by retry config)
        status = getattr(getattr(exc, "response", None), "status_code", None)
        if status and 400 <= status < 500 and status != 429:
            raise ParsingError(f"LlamaParse returned client error {status}: {exc}") from exc
        raise ParsingError(f"LlamaParse failed: {exc}") from exc
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)

    # LlamaParse markdown doesn't always preserve page boundaries as explicit
    # markers. We try to detect "---" page-break conventions or treat each
    # document section as a block at page index 1 (acceptable for tables/figures
    # where cross-page citation is less critical than prose).
    blocks: list[PageBlock] = []
    for doc in documents:
        pages = doc.text.split("\f")  # form-feed is LlamaParse's page separator
        if len(pages) > 1:
            for page_num, page_text in enumerate(pages, start=1):
                text = page_text.strip()
                if text:
                    blocks.append((page_num, text))
        else:
            # Fallback: whole doc as page 1
            text = doc.text.strip()
            if text:
                blocks.append((1, text))

    logger.info("LlamaParse: extracted %d page blocks.", len(blocks))
    return blocks


# --------------------------------------------------------------------------- #
#  Auto-routing heuristic                                                      #
# --------------------------------------------------------------------------- #

def _detect_hint(content: bytes) -> DocumentHint:
    """Inspect PDF byte content and suggest a routing hint."""
    try:
        import fitz
        doc = fitz.open(stream=content, filetype="pdf")
        image_count = sum(len(page.get_images()) for page in doc)
        text_len = sum(len(page.get_text("text")) for page in doc)
        page_count = len(doc)
        doc.close()

        images_per_page = image_count / max(page_count, 1)
        chars_per_page = text_len / max(page_count, 1)

        if chars_per_page < 100:
            return "scanned"
        if images_per_page > 2:
            return "tables"
        return "mostly_text"
    except Exception:
        return "mostly_text"  # safe fallback


# --------------------------------------------------------------------------- #
#  Public dispatcher                                                           #
# --------------------------------------------------------------------------- #

async def parse_pdf(
    content: bytes,
    hint: DocumentHint = "auto",
    filename: str = "document.pdf",
) -> tuple[list[PageBlock], bool]:
    """Dispatch PDF parsing to the appropriate backend based on hint.

    Returns a tuple: (list of (page_number, text) tuples, is_markdown bool)
    Raises ParsingError on failure.
    """
    logger.info("parse_pdf: file=%s hint=%s size=%d bytes", filename, hint, len(content))

    if hint == "auto":
        hint = _detect_hint(content)
        logger.info("Auto-detected routing hint: %s", hint)

    if hint == "mostly_text":
        loop = asyncio.get_running_loop()
        blocks = await loop.run_in_executor(None, _parse_pdf_pymupdf, content)
        return blocks, False
    else:
        # tables / equations / scanned → LlamaParse
        blocks = await _parse_pdf_llamaparse(content, filename=filename)
        return blocks, True


async def parse_pptx(content: bytes) -> list[PageBlock]:
    """Convert PPTX bytes to (slide_number, text) tuples."""
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(content))
        blocks: list[PageBlock] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                blocks.append((slide_num, "\n".join(texts)))

        logger.info("PPTX: extracted %d slide blocks.", len(blocks))
        return blocks
    except Exception as exc:
        raise ParsingError(f"PPTX parsing failed: {exc}") from exc


async def parse_docx(content: bytes) -> list[PageBlock]:
    """Convert DOCX bytes to a single page block (DOCX has no page metadata)."""
    from docx import Document

    _HEADING_MARKERS = {"Heading 1": "#", "Heading 2": "##", "Heading 3": "###"}

    try:
        doc = Document(io.BytesIO(content))
        lines: list[str] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            for heading_style, marker in _HEADING_MARKERS.items():
                if heading_style in para.style.name:
                    text = f"{marker} {text}"
                    break
            lines.append(text)

        logger.info("DOCX: extracted %d paragraphs.", len(lines))
        return [(1, "\n\n".join(lines))] if lines else []
    except Exception as exc:
        raise ParsingError(f"DOCX parsing failed: {exc}") from exc


async def parse_file(
    filename: str,
    content_type: str,
    content: bytes,
    hint: DocumentHint = "auto",
) -> tuple[list[PageBlock], bool]:
    """Dispatch to the correct parser based on MIME type.

    Returns a tuple: (list of (page_number, text) tuples, is_markdown bool)
    Raises ParsingError for unsupported types or parse failures.
    """
    logger.debug("Dispatching file '%s' (content_type=%s, size=%d)", filename, content_type, len(content))

    if content_type == "application/pdf":
        return await parse_pdf(content, hint=hint, filename=filename)

    if "pptx" in content_type or "powerpoint" in content_type or "presentationml" in content_type:
        blocks = await parse_pptx(content)
        return blocks, False

    if "docx" in content_type or "wordprocessingml" in content_type:
        blocks = await parse_docx(content)
        return blocks, False

    if content_type == "text/plain":
        text = content.decode("utf-8", errors="ignore")
        blocks = [(1, text)] if text.strip() else []
        return blocks, False

    raise ParsingError(
        f"Unsupported content type '{content_type}' for file '{filename}'. "
        "Accepted: application/pdf, pptx, docx, text/plain."
    )
